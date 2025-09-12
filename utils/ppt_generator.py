import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class PPTGenerator:
    def __init__(self, output_dir="generated_ppts"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def generate_ppt(self, final_json: dict, file_name: str = "presentation.pptx"):
        """
        Convert the final formatted JSON into a PowerPoint file.
        """
        try:
            prs = Presentation()

            # Set slide layout: Title + Content
            slide_layout = prs.slide_layouts[1]

            for slide_data in final_json["slides"]:
                slide_title = slide_data["title"]
                content_points = slide_data["content"]

                # Create a new slide
                slide = prs.slides.add_slide(slide_layout)

                # Set slide title
                title_placeholder = slide.shapes.title
                title_placeholder.text = slide_title
                title_placeholder.text_frame.paragraphs[0].font.size = Pt(28)
                title_placeholder.text_frame.paragraphs[0].font.bold = True
                title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(34, 34, 34)

                # Add content points
                body_placeholder = slide.placeholders[1]
                tf = body_placeholder.text_frame
                tf.clear()  # Clear default text

                for point in content_points:
                    statement = point.get("statement", " No statement provided")
                    status = point.get("status", "needs_review")
                    citation = point.get("citation", "Not verified")
                    design_hint = point.get("design_hint", "")

                    # Add main statement
                    p = tf.add_paragraph()
                    p.text = f"• {statement}"
                    p.font.size = Pt(18)
                    p.font.bold = status == "needs_review"  # Highlight review-needed points
                    p.font.color.rgb = RGBColor(0, 102, 204) if status == "accurate" else RGBColor(204, 51, 51)
                    p.alignment = PP_ALIGN.LEFT

                    # Optionally, add citations in smaller font
                    if citation and citation != "Not verified":
                        citation_p = tf.add_paragraph()
                        citation_p.text = f"    (Citation: {citation})"
                        citation_p.font.size = Pt(12)
                        citation_p.font.color.rgb = RGBColor(100, 100, 100)
                        citation_p.alignment = PP_ALIGN.LEFT

                    # Optionally, include design hints in lighter text
                    if design_hint:
                        hint_p = tf.add_paragraph()
                        hint_p.text = f"    [Design hint: {design_hint}]"
                        hint_p.font.size = Pt(11)
                        hint_p.font.color.rgb = RGBColor(120, 120, 120)
                        hint_p.alignment = PP_ALIGN.LEFT

            # Save the PowerPoint file
            file_path = os.path.join(self.output_dir, file_name)
            prs.save(file_path)
            return file_path

        except Exception as e:
            raise RuntimeError(f" Failed to generate PPT: {e}")
